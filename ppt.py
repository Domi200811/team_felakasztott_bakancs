from pptx import Presentation
from pptx.util import Inches, Pt

def create_hangman_presentation():
    # Prezentáció objektum létrehozása
    prs = Presentation()

    # --- 1. Dia: Címlap ---
    slide_layout = prs.slide_layouts[0] # 0 = Címlap elrendezés
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "HANGMAN"
    subtitle.text = "Készítette:\nJuhász Bálint, Rónai Márton János, Pék Dominik"

    # --- 2. Dia: A játékról ---
    slide_layout = prs.slide_layouts[1] # 1 = Cím és tartalom elrendezés
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "A játékról"
    tf = slide.placeholders[1].text_frame
    
    tf.text = "A program az akasztófa (hangman) játék Pythonban megvalósítva."
    tf.add_paragraph().text = "A felhasználónak ki kell találnia egy véletlenszerű országot vagy fővárost."
    tf.add_paragraph().text = "Összesen 6 esély (hiba) megengedett."
    tf.add_paragraph().text = "Amennyiben ez nem sikerül, a játékost felakasztják."

    # --- 3. Dia: Működési elv ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Működési elv"
    tf = slide.placeholders[1].text_frame
    
    tf.text = "Külső fájl forrásból (txt/json) beolvassuk az országokat és fővárosokat."
    tf.add_paragraph().text = "A felhasználó eldönti, mit szeretne kitalálni (1: Ország, 2: Főváros)."
    tf.add_paragraph().text = "A program halmazokat (set) használ a már tippelt betűk tárolására, így egy betűt nem lehet kétszer beírni."
    tf.add_paragraph().text = "A mask_word függvény a kitalálandó szót rejtett formátumba (-) konvertálja, de a szóközöket megtartja."

    # --- 4. Dia: Újdonság: Toplista (Scoreboard) ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Újdonság: Toplista (Scoreboard)"
    tf = slide.placeholders[1].text_frame
    
    tf.text = "Adatok mentése és betöltése egy 'scoreboard.json' fájlba."
    tf.add_paragraph().text = "Ha a játékos nyer, a program kiszámolja, hány tippre volt szüksége."
    tf.add_paragraph().text = "A listát automatikusan pontszám szerint rendezi (a legkevesebb tipp a legjobb)."
    tf.add_paragraph().text = "Csak a Top 10 legjobb játékos marad a dicsőségfalon!"

    # --- 5. Dia: Kihívások és Nehézségek (Mémek helye) ---
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # 5 = Csak cím
    slide.shapes.title.text = "Kihívások és Nehézségek a fejlesztés során"
    
    # Helyőrző szövegdoboz a mémeknek
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "In memory of Marci...\n[IDE JÖHETNEK A MÉMEK: Csontváz, Stewie, Hazmat suit]\nIde szúrjátok be a PPT-ben a vicces képeket, és meséljétek el, melyik kódrészlettel (pl. JSON vagy a szóközök kezelése) szenvedtetek a legtöbbet!"
    p.font.italic = True

    # --- 6. Dia: Demo ---
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Élő Demo"
    tf = slide.placeholders[1].text_frame
    
    tf.text = "Próbáljuk ki élőben!"
    tf.add_paragraph().text = "Kérek egy tippet a közönségből!"
    tf.add_paragraph().text = "Link a kódhoz (OneCompiler): https://onecompiler.com/python/44ckz43bk"

    # Prezentáció mentése
    prs.save('Hangman_Prezentacio_Uj.pptx')
    print("A 'Hangman_Prezentacio_Uj.pptx' sikeresen legenerálva!")

if __name__ == "__main__":
    create_hangman_presentation()